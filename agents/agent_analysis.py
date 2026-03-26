"""
Agent 5: Conversational Policy Analysis Agent
Uses ONLY gpt-4o-mini via chat.completions — no Responses API, no search models.
"""
from openai import OpenAI
import json
import sqlite3
import os
import pandas as pd
import streamlit as st
from datetime import date, datetime

PROJECT_END = date(2028, 6, 30)
TODAY       = date.today()
DB_PATH     = "/tmp/agent5_conversations.db"   # /tmp is writable on Streamlit Cloud

VERDICTS = ["CRITICAL", "UPGRADE NOW", "EXTEND + PLAN", "REPLACE", "CLOUD MIGRATE", "MONITOR"]

# ── OS Family Categorization ─────────────────────────────────────────────────
OS_FAMILY_RULES = [
    ("Windows Family",      ["Windows"],       "Desktop 10/11 & Server 2012–2025"),
    ("RHEL/Clone Family",   ["RHEL", "Red Hat", "AlmaLinux", "Alma", "Rocky", "Oracle Linux",
                             "CentOS", "CentOS Stream"],
                                               "RHEL, Alma, Rocky, Oracle Linux, CentOS"),
    ("Debian/Ubuntu Family", ["Ubuntu", "Debian"], "Ubuntu 14.04+ & Debian 8–13"),
    ("SUSE Family",         ["SLES", "SUSE", "openSUSE"], "SLES 11–16"),
    ("Legacy Unix",         ["AIX", "HP-UX", "Solaris", "Tru64"], "AIX, HP-UX, Solaris, Tru64"),
    ("BSD & VMS",           ["FreeBSD", "OpenBSD", "NetBSD", "OpenVMS"], "FreeBSD & OpenVMS"),
    ("Apple",               ["macOS", "iOS", "iPadOS"], "macOS 13–26"),
]

def categorize_os_families(os_df):
    """Categorize the OS dataframe into families. Returns dict: family_name → list of OS versions."""
    families = {}
    uncategorized = []
    for _, row in os_df.iterrows():
        ver = str(row.get("OS Version", ""))
        matched = False
        for fam_name, keywords, _ in OS_FAMILY_RULES:
            if any(kw.lower() in ver.lower() for kw in keywords):
                families.setdefault(fam_name, []).append(ver)
                matched = True
                break
        if not matched and ver:
            uncategorized.append(ver)
    # Add uncategorized to "Other" if any
    if uncategorized:
        families["Other"] = uncategorized
    return families

def get_family_display():
    """Return list of (family_name, description, emoji) for UI display."""
    return [
        ("Windows Family",       "Desktop 10/11 & Server 2012–2025", "🪟"),
        ("RHEL/Clone Family",    "RHEL, Alma, Rocky, Oracle Linux, CentOS", "🐧"),
        ("Debian/Ubuntu Family", "Ubuntu 14.04+ & Debian 8–13", "🐧"),
        ("SUSE Family",          "SLES 11–16", "🦎"),
        ("Legacy Unix",          "AIX, HP-UX, Solaris, Tru64", "🖥️"),
        ("BSD & VMS",            "FreeBSD & OpenVMS", "👾"),
        ("Apple",                "macOS 13–26", "🍎"),
        ("Other",                "Other OS not listed above", "❓"),
    ]

# ── OS Migration Guiding Principles (default knowledge base) ─────────────────
DEFAULT_OS_PRINCIPLES = {
    "Windows Family": {
        "cloud_map": {"Azure": "Azure / Azure Gov", "AWS": "AWS", "GCP": "GCP", "OCI": "OCI",
                       "GovCloud": "Azure Gov / AWS GovCloud", "_default": "Azure / Azure Gov"},
        "upgrade": "COTS: Upgrade to Server 2022/2025 to maintain vendor support and security compliance.",
        "replace":  "Custom: Replace Windows OS with Linux Containers (.NET Core) to eliminate license costs.",
    },
    "RHEL/Clone Family": {
        "cloud_map": {"AWS": "AWS / GovCloud", "Azure": "Azure", "GCP": "GCP", "OCI": "OCI",
                       "GovCloud": "AWS GovCloud", "_default": "AWS / GovCloud"},
        "upgrade": "COTS: Upgrade to RHEL 9.x or Oracle Linux 9.x as required by the software vendor's certification.",
        "replace":  "Custom: Replace legacy Linux with Amazon Linux 2023 to leverage AWS-native security and performance.",
    },
    "Debian/Ubuntu Family": {
        "cloud_map": {"GCP": "GCP / GDC", "AWS": "AWS", "Azure": "Azure", "OCI": "OCI",
                       "GovCloud": "GCP Assured Workloads", "_default": "GCP / GDC"},
        "upgrade": "COTS: Upgrade to the latest LTS (Long Term Support) version to ensure a stable 5-year security patch window.",
        "replace":  "Custom: Replace the standalone OS with Google Container-Optimized OS for microservices workloads.",
    },
    "SUSE Family": {
        "cloud_map": {"Azure": "Azure / AWS", "AWS": "Azure / AWS", "GCP": "GCP", "OCI": "OCI",
                       "GovCloud": "Azure Gov", "_default": "Azure / AWS"},
        "upgrade": "COTS: Upgrade to SLES 15 SP5+ specifically for SAP HANA or high-availability enterprise applications.",
        "replace":  "Custom: Replace manual SLES installs with Automated Cloud Images managed via SUSE Manager.",
    },
    "Legacy Unix": {
        "cloud_map": {"_default": "On-Prem / Bare Metal"},
        "upgrade": "None: No cloud upgrade path exists. Maintain \"As-Is\" on existing hardware with network isolation.",
        "replace":  "COTS/Custom: Mandatory Replacement. Transition the workload to Modern x86 Linux on the preferred cloud provider.",
    },
    "BSD & VMS": {
        "cloud_map": {"_default": "Specialized / Hybrid"},
        "upgrade": "COTS: Upgrade to the latest stable release branch supported by the hardware abstraction layer.",
        "replace":  "Custom: Replace with Hardened Linux Distributions to consolidate the technical skill set required for maintenance.",
    },
    "Apple": {
        "cloud_map": {"AWS": "AWS / Local", "Azure": "Azure", "GCP": "GCP",
                       "_default": "AWS / Local"},
        "upgrade": "Custom: Upgrade to the latest N-1 Version to ensure compatibility with modern developer toolchains (Xcode).",
        "replace":  "Custom: Replace physical on-prem hardware with Managed EC2 Mac Instances for cloud-based CI/CD.",
    },
}

# ── DB Migration Principles ───────────────────────────────────────────────────
DEFAULT_DB_PRINCIPLES = {
    "SQL Server":     {"upgrade": "Upgrade to SQL Server 2022 for enhanced security and Azure Hybrid Benefit.", "replace": "Migrate to Azure SQL Managed Instance or AWS RDS for managed operations."},
    "Oracle":         {"upgrade": "Upgrade to Oracle 19c/23ai for extended premier support.", "replace": "Migrate to PostgreSQL or AWS Aurora PostgreSQL to eliminate licensing costs."},
    "PostgreSQL":     {"upgrade": "Upgrade to PostgreSQL 16/17 for performance and security.", "replace": "Move to cloud-managed (RDS/Cloud SQL/Azure Flexible) for zero-maintenance operations."},
    "MySQL":          {"upgrade": "Upgrade to MySQL 8.4 LTS for long-term security patches.", "replace": "Migrate to Aurora MySQL or Cloud SQL for auto-scaling and managed backups."},
    "MongoDB":        {"upgrade": "Upgrade to MongoDB 7.x/8.x for improved aggregation and security.", "replace": "Move to MongoDB Atlas for fully managed cloud-native operations."},
    "IBM Db2":        {"upgrade": "Upgrade to Db2 11.5 with latest fix packs for extended support.", "replace": "Migrate to PostgreSQL on cloud for cost reduction and talent availability."},
    "MariaDB":        {"upgrade": "Upgrade to MariaDB 11.4 LTS for enterprise stability.", "replace": "Consolidate onto Aurora MySQL or Cloud SQL MySQL-compatible for managed ops."},
}

# ── Web Server Migration Principles ──────────────────────────────────────────
DEFAULT_WS_PRINCIPLES = {
    "IIS":     {"upgrade": "Upgrade to IIS 10.0 on Windows Server 2022/2025 for HTTP/3 and TLS 1.3.", "replace": "Replace with Nginx or cloud ALB/CDN for static content and reverse proxy."},
    "Nginx":   {"upgrade": "Upgrade to Nginx 1.26+ stable for latest security patches and HTTP/3.", "replace": "Replace with cloud-native load balancers (ALB/Cloud Load Balancing) for auto-scaling."},
    "Apache":  {"upgrade": "Upgrade to Apache 2.4.62+ for security patches and mod_http2.", "replace": "Replace with Nginx or cloud-managed API Gateway for modern workloads."},
}

# ── App Server Migration Principles ──────────────────────────────────────────
DEFAULT_AS_PRINCIPLES = {
    "Tomcat":          {"upgrade": "Upgrade to Tomcat 10.1.x (Jakarta EE 10) for modern servlet support.", "replace": "Containerize with Docker/K8s and deploy on cloud-managed EKS/GKE/AKS."},
    "JBoss/WildFly":   {"upgrade": "Upgrade to JBoss EAP 8.0 or WildFly 33 for Jakarta EE 10.", "replace": "Migrate to Quarkus on Kubernetes for cloud-native microservices."},
    "WebSphere":       {"upgrade": "Migrate to WebSphere Liberty for lightweight, cloud-ready deployment.", "replace": "Re-platform to Spring Boot or Quarkus on Kubernetes."},
    "Kafka":           {"upgrade": "Upgrade to Kafka 3.x (KRaft) to eliminate ZooKeeper dependency.", "replace": "Move to managed Kafka (Amazon MSK / Confluent Cloud) for zero-ops."},
    "RabbitMQ":        {"upgrade": "Upgrade to RabbitMQ 3.13.x for latest security and performance.", "replace": "Migrate to Amazon MQ or CloudAMQP for managed message brokering."},
    "Kubernetes":      {"upgrade": "Upgrade to K8s 1.31+ and maintain N-2 version policy.", "replace": "Adopt managed K8s (EKS/GKE/AKS) to eliminate control plane management overhead."},
}

# ── Framework Migration Principles ───────────────────────────────────────────
DEFAULT_FW_PRINCIPLES = {
    ".NET":            {"upgrade": "Upgrade to .NET 8 LTS for cross-platform support and performance.", "replace": "Migrate legacy .NET Framework apps to .NET 8 containers on Linux."},
    "Spring Boot":     {"upgrade": "Upgrade to Spring Boot 3.4.x for Jakarta EE 10 and virtual threads.", "replace": "Refactor monoliths to Spring Cloud microservices on Kubernetes."},
    "PHP":             {"upgrade": "Upgrade to PHP 8.3/8.4 for JIT compilation and type safety.", "replace": "Containerize PHP-FPM workloads on Kubernetes with auto-scaling."},
    "React":           {"upgrade": "Upgrade to React 19 for Server Components and Suspense improvements.", "replace": "Deploy on cloud CDN (CloudFront/Cloud CDN) with SSR via serverless functions."},
    "Angular":         {"upgrade": "Upgrade to Angular 18/19 for signals and improved hydration.", "replace": "Serve via containerized SSR on cloud Kubernetes with edge caching."},
    "Node.js":         {"upgrade": "Upgrade to Node.js 22 LTS for V8 improvements and security.", "replace": "Deploy on serverless (Lambda/Cloud Functions) for event-driven workloads."},
    "Java":            {"upgrade": "Upgrade to Java 21 LTS for virtual threads and pattern matching.", "replace": "Adopt GraalVM native images for cloud-optimized startup and memory."},
    "Python":          {"upgrade": "Upgrade to Python 3.12/3.13 for performance and security patches.", "replace": "Deploy on serverless (Lambda/Cloud Functions) or containerized ML pipelines."},
    "Django":          {"upgrade": "Upgrade to Django 5.2 LTS for async views and composite PKs.", "replace": "Containerize Django apps on Kubernetes with managed PostgreSQL backend."},
    "Vue.js":          {"upgrade": "Upgrade to Vue.js 3.x with Composition API for maintainability.", "replace": "Deploy on cloud CDN with SSR via Nuxt.js on serverless."},
    "Ruby on Rails":   {"upgrade": "Upgrade to Rails 8.0 for Hotwire-first architecture.", "replace": "Containerize on Kubernetes with managed PostgreSQL/Redis backends."},
}

PRINCIPLES_TABLE_SYSTEM = f"""You are Agent 5 — a senior IT migration strategist at Infosys.
Given the user's technology landscape and preferred cloud, generate a deep-dive guiding
principles table covering OS, Database, Web Server, App Server, and Framework categories.

For each item, provide:
- Category (OS/DB/Web Server/App Server/Framework)
- Technology name
- Cloud target (adapted to user's preference)
- Upgrade Principle (COTS/vendor-supported path)
- Replacement Principle (modernization/cloud-native path)

Be specific, actionable, and reference actual product versions and cloud services.
Return ONLY a JSON array:
[{{"category": "...", "technology": "...", "cloud_target": "...", "upgrade_principle": "...", "replacement_principle": "..."}}]

Today: {TODAY}. Project ends 30 Jun 2028."""


def generate_principles_table(selected_families, cloud_name, cloud_key,
                               agent=None, db_df=None, ws_df=None, as_df=None, fw_df=None):
    """
    Generate the deep-dive guiding principles table covering OS, DB, WS, AS, FW.
    Returns list of dicts: category, technology, cloud_target, upgrade_principle, replacement_principle.
    """
    rows = []

    # ── OS Families ──────────────────────────────────────────────────────────
    for fam in selected_families:
        if fam == "Other":
            continue
        defaults = DEFAULT_OS_PRINCIPLES.get(fam)
        if defaults:
            cloud_map = defaults.get("cloud_map", {})
            cloud_target = cloud_map.get(cloud_key, cloud_map.get("_default", cloud_name))
            rows.append({
                "category": "OS",
                "technology": fam,
                "cloud_target": cloud_target,
                "upgrade_principle": defaults["upgrade"],
                "replacement_principle": defaults["replace"],
            })

    # ── Databases (from live data) ───────────────────────────────────────────
    if db_df is not None and not db_df.empty:
        db_names = db_df["Database"].dropna().unique().tolist()
        for db_name in sorted(set(db_names)):
            matched = None
            for key in DEFAULT_DB_PRINCIPLES:
                if key.lower() in db_name.lower():
                    matched = DEFAULT_DB_PRINCIPLES[key]
                    break
            if matched:
                rows.append({
                    "category": "Database",
                    "technology": db_name,
                    "cloud_target": cloud_name,
                    "upgrade_principle": matched["upgrade"],
                    "replacement_principle": matched["replace"],
                })

    # ── Web Servers (from live data) ─────────────────────────────────────────
    if ws_df is not None and not ws_df.empty:
        ws_names = ws_df["Web Server"].dropna().unique().tolist()
        for ws_name in sorted(set(ws_names)):
            matched = None
            for key in DEFAULT_WS_PRINCIPLES:
                if key.lower() in ws_name.lower():
                    matched = DEFAULT_WS_PRINCIPLES[key]
                    break
            if matched:
                rows.append({
                    "category": "Web Server",
                    "technology": ws_name,
                    "cloud_target": cloud_name,
                    "upgrade_principle": matched["upgrade"],
                    "replacement_principle": matched["replace"],
                })

    # ── App Servers (from live data) ─────────────────────────────────────────
    if as_df is not None and not as_df.empty:
        as_names = as_df["App Server"].dropna().unique().tolist()
        for as_name in sorted(set(as_names)):
            matched = None
            for key in DEFAULT_AS_PRINCIPLES:
                if key.lower() in as_name.lower():
                    matched = DEFAULT_AS_PRINCIPLES[key]
                    break
            if matched:
                rows.append({
                    "category": "App Server",
                    "technology": as_name,
                    "cloud_target": cloud_name,
                    "upgrade_principle": matched["upgrade"],
                    "replacement_principle": matched["replace"],
                })

    # ── Frameworks (from live data) ──────────────────────────────────────────
    if fw_df is not None and not fw_df.empty:
        fw_names = fw_df["Framework"].dropna().unique().tolist()
        for fw_name in sorted(set(fw_names)):
            matched = None
            for key in DEFAULT_FW_PRINCIPLES:
                if key.lower() in fw_name.lower():
                    matched = DEFAULT_FW_PRINCIPLES[key]
                    break
            if matched:
                rows.append({
                    "category": "Framework",
                    "technology": fw_name,
                    "cloud_target": cloud_name,
                    "upgrade_principle": matched["upgrade"],
                    "replacement_principle": matched["replace"],
                })

    # ── AI Enhancement ───────────────────────────────────────────────────────
    if agent and rows:
        try:
            prompt = (
                f"User's cloud preference: {cloud_name}\n"
                f"User's OS families: {', '.join(selected_families)}\n\n"
                f"Here are the default principles across all categories:\n"
                + json.dumps(rows, indent=2) +
                f"\n\nEnhance these to be tailored to {cloud_name}. "
                f"Reference specific {cloud_name} services. Keep the same JSON structure. "
                f"Return ONLY the JSON array."
            )
            resp = agent.client.chat.completions.create(
                model=agent.model, max_tokens=4000,
                messages=[
                    {"role": "system", "content": PRINCIPLES_TABLE_SYSTEM},
                    {"role": "user", "content": prompt}
                ])
            text = resp.choices[0].message.content.strip()
            if "```" in text:
                text = text.split("```json")[-1].split("```")[0] if "```json" in text \
                       else text.split("```")[1].split("```")[0]
            s, e = text.find("["), text.rfind("]")
            if s != -1 and e > s:
                enhanced = json.loads(text[s:e+1])
                if enhanced and len(enhanced) >= len(rows) // 2:
                    rows = enhanced
        except Exception:
            pass

    return rows


# =============================================================================
# FEATURE: Migration Wave Planner
# =============================================================================
def assign_migration_waves(table_data, os_df=None, db_df=None, ws_df=None, as_df=None, fw_df=None):
    """Assign each technology to a migration wave based on risk scores and EOL dates."""
    waves = []
    # Build risk lookup from all dataframes
    risk_lookup = {}
    for df, name_col in [
        (os_df, "OS Version"), (db_df, "Database"), (ws_df, "Web Server"),
        (as_df, "App Server"), (fw_df, "Framework")
    ]:
        if df is None or df.empty or name_col not in df.columns:
            continue
        for _, row in df.iterrows():
            key = str(row.get(name_col, ""))
            risk_lookup[key] = {
                "risk_score": row.get("Risk Score", 50),
                "risk_level": row.get("Risk Level", "MEDIUM"),
                "days_eol": row.get("Days Until EOL", 999),
                "status": row.get("Status", "Supported"),
            }

    for item in table_data:
        tech = item.get("technology", item.get("os_family", ""))
        cat = item.get("category", "OS")

        # Find best match in risk lookup
        best_risk = None
        for key, data in risk_lookup.items():
            if tech.lower() in key.lower() or key.lower() in tech.lower():
                if best_risk is None or data["risk_score"] > best_risk["risk_score"]:
                    best_risk = data

        if best_risk is None:
            best_risk = {"risk_score": 50, "risk_level": "MEDIUM", "days_eol": 999, "status": "Supported"}

        # Assign wave
        score = best_risk["risk_score"]
        days = best_risk["days_eol"] if best_risk["days_eol"] is not None else 999
        status = str(best_risk.get("status", "")).lower()

        if score >= 75 or days < 0 or status == "end of life":
            wave = 1; wave_name = "Wave 1: Critical EOL"; timeline = "0–6 months"; urgency = "CRITICAL"
        elif score >= 50 or (0 <= days <= 365) or status == "expiring soon":
            wave = 2; wave_name = "Wave 2: High Risk"; timeline = "6–12 months"; urgency = "HIGH"
        elif score >= 30 or (365 < days <= 730):
            wave = 3; wave_name = "Wave 3: Plan Ahead"; timeline = "12–18 months"; urgency = "MEDIUM"
        else:
            wave = 4; wave_name = "Wave 4: Monitor"; timeline = "18–24 months"; urgency = "LOW"

        waves.append({
            **item,
            "wave": wave, "wave_name": wave_name,
            "timeline": timeline, "urgency": urgency,
            "risk_score": best_risk["risk_score"],
            "days_eol": days,
        })

    return sorted(waves, key=lambda x: (x["wave"], -x["risk_score"]))


# =============================================================================
# FEATURE: Dependency Mapping
# =============================================================================
COMMON_DEPENDENCIES = [
    # (app_server/framework, typical_os, typical_db, note)
    ("Tomcat", "RHEL/Clone Family", "PostgreSQL", "Java servlet container commonly on RHEL with PostgreSQL/MySQL"),
    ("Tomcat", "Windows Family", "SQL Server", "Java on Windows typically backed by SQL Server"),
    ("JBoss/WildFly", "RHEL/Clone Family", "Oracle", "Enterprise Java on RHEL often paired with Oracle DB"),
    ("WebSphere", "RHEL/Clone Family", "IBM Db2", "IBM stack: WAS + RHEL + Db2"),
    ("WebSphere", "Windows Family", "SQL Server", "WebSphere on Windows with SQL Server"),
    ("IIS", "Windows Family", "SQL Server", ".NET/IIS stack requires Windows + SQL Server"),
    (".NET", "Windows Family", "SQL Server", ".NET Framework requires Windows; .NET 8 can run on Linux"),
    ("Spring Boot", "RHEL/Clone Family", "PostgreSQL", "Spring Boot commonly deployed on Linux with PostgreSQL"),
    ("Django", "Debian/Ubuntu Family", "PostgreSQL", "Django + Ubuntu + PostgreSQL is the standard Python stack"),
    ("PHP", "Debian/Ubuntu Family", "MySQL", "LAMP/LEMP stack: Linux + Nginx/Apache + MySQL + PHP"),
    ("Node.js", "Debian/Ubuntu Family", "MongoDB", "MERN/MEAN stack: Node.js + MongoDB on Ubuntu"),
    ("Kafka", "RHEL/Clone Family", None, "Kafka runs on Linux; ZooKeeper dependency removed in KRaft mode"),
    ("Kubernetes", "Debian/Ubuntu Family", None, "K8s nodes typically run Ubuntu or container-optimized OS"),
    ("Nginx", "Debian/Ubuntu Family", None, "Nginx commonly deployed on Ubuntu/Debian"),
    ("Apache", "RHEL/Clone Family", None, "Apache httpd commonly on RHEL-based systems"),
    ("React", None, None, "Frontend framework — deployed via CDN, no OS/DB dependency at runtime"),
    ("Angular", None, None, "Frontend framework — deployed via CDN, no OS/DB dependency at runtime"),
    ("Vue.js", None, None, "Frontend framework — deployed via CDN, no OS/DB dependency at runtime"),
]

def generate_dependency_map(table_data):
    """Generate dependency chains showing which technologies must move together."""
    deps = []
    techs = {r.get("technology", r.get("os_family", "")).lower(): r for r in table_data}

    for app, os_fam, db, note in COMMON_DEPENDENCIES:
        app_match = None
        for key, item in techs.items():
            if app.lower() in key.lower():
                app_match = item
                break
        if not app_match:
            continue

        chain = {"source": app_match.get("technology", app),
                 "source_category": app_match.get("category", ""),
                 "depends_on": [], "note": note}

        if os_fam:
            for key, item in techs.items():
                if os_fam.lower().split("/")[0].lower() in key.lower() or key in os_fam.lower():
                    chain["depends_on"].append({
                        "technology": item.get("technology", os_fam),
                        "category": "OS", "type": "runs_on"})
                    break

        if db:
            for key, item in techs.items():
                if db.lower() in key.lower():
                    chain["depends_on"].append({
                        "technology": item.get("technology", db),
                        "category": "Database", "type": "uses_db"})
                    break

        if chain["depends_on"]:
            deps.append(chain)

    return deps


# =============================================================================
# FEATURE: Cost Estimator
# =============================================================================
COST_ESTIMATES = {
    # technology_keyword: (upgrade_cost_per_unit, replace_cost_per_unit, do_nothing_annual, unit)
    "Windows": ("$2K–5K/server (in-place upgrade)", "$15K–30K/server (re-platform to Linux)", "$200–400/server/yr ESU", "server"),
    "RHEL": ("$500–1K/server (minor version)", "$3K–8K/server (distro migration)", "$800–1.3K/server/yr subscription", "server"),
    "Ubuntu": ("$200–500/server (LTS upgrade)", "$3K–8K/server (distro migration)", "$500/server/yr Ubuntu Pro", "server"),
    "SUSE": ("$500–1K/server (SP upgrade)", "$5K–10K/server (migration)", "$1K–2K/server/yr subscription", "server"),
    "Legacy Unix": ("N/A (no upgrade path)", "$50K–200K/system (re-platform)", "$10K–50K/system/yr maintenance", "system"),
    "BSD": ("$200–500/server", "$5K–15K/server (migration)", "$1K–3K/server/yr", "server"),
    "Apple": ("$0 (free OS upgrade)", "$3K–5K/device (EC2 Mac)", "$0 (included with hardware)", "device"),
    "SQL Server": ("$5K–15K/instance (version upgrade)", "$20K–50K/instance (cloud migration)", "$1.4K/core/yr ESU", "instance"),
    "Oracle": ("$10K–30K/instance (patch set)", "$30K–100K/instance (PostgreSQL migration)", "22% of license/yr support", "instance"),
    "PostgreSQL": ("$1K–3K/instance (version upgrade)", "$5K–15K/instance (managed cloud)", "$0 (community) or $2K/yr", "instance"),
    "MySQL": ("$1K–3K/instance", "$5K–15K/instance (Aurora)", "$2K–5K/yr Enterprise", "instance"),
    "MongoDB": ("$2K–5K/instance", "$5K–20K/instance (Atlas)", "$3K–10K/yr Enterprise", "instance"),
    "Db2": ("$5K–15K/instance", "$30K–80K/instance (PostgreSQL)", "$15K–30K/yr maintenance", "instance"),
    "IIS": ("$0 (included with Windows)", "$5K–15K (Nginx/cloud ALB)", "$0 (Windows license)", "server"),
    "Nginx": ("$0 (free upgrade)", "$3K–8K (cloud ALB)", "$0 (OSS) or $2.5K/yr Plus", "server"),
    "Apache": ("$0 (free upgrade)", "$3K–8K (cloud ALB/Nginx)", "$0 (OSS)", "server"),
    "Tomcat": ("$1K–3K/instance (version)", "$5K–15K/instance (containerize)", "$0 (OSS)", "instance"),
    "JBoss": ("$3K–8K/instance", "$10K–25K/instance (Quarkus)", "$5K–15K/yr subscription", "instance"),
    "WebSphere": ("$5K–15K/instance (Liberty)", "$15K–40K/instance (re-platform)", "$10K–25K/yr license", "instance"),
    "Kafka": ("$2K–5K/cluster", "$5K–20K (managed MSK/Confluent)", "$0 (OSS) or $10K/yr", "cluster"),
    "RabbitMQ": ("$1K–3K/cluster", "$3K–10K (managed)", "$0 (OSS)", "cluster"),
    "Kubernetes": ("$1K–3K/cluster (version)", "$5K–15K (managed EKS/GKE)", "$0 (OSS)", "cluster"),
    ".NET": ("$2K–8K/app (framework upgrade)", "$10K–40K/app (containerize)", "$0–5K/yr", "application"),
    "Spring Boot": ("$2K–5K/app (version upgrade)", "$5K–20K/app (microservices)", "$0 (OSS)", "application"),
    "PHP": ("$1K–3K/app", "$5K–15K/app (containerize)", "$0 (OSS)", "application"),
    "React": ("$1K–3K/app (version upgrade)", "$2K–8K/app (SSR/CDN)", "$0 (OSS)", "application"),
    "Angular": ("$1K–3K/app", "$3K–10K/app (rewrite)", "$0 (OSS)", "application"),
    "Node.js": ("$1K–3K/app", "$3K–10K/app (serverless)", "$0 (OSS)", "application"),
    "Java": ("$1K–3K/app (LTS upgrade)", "$5K–15K/app (GraalVM native)", "$0 (OpenJDK)", "application"),
    "Python": ("$500–2K/app", "$3K–10K/app (containerize)", "$0 (OSS)", "application"),
    "Django": ("$1K–3K/app", "$5K–15K/app (containerize)", "$0 (OSS)", "application"),
    "Vue.js": ("$1K–2K/app", "$2K–5K/app (Nuxt SSR)", "$0 (OSS)", "application"),
    "Rails": ("$2K–5K/app", "$5K–20K/app (containerize)", "$0 (OSS)", "application"),
}

COST_AI_SYSTEM = f"""You are a cloud migration cost analyst with expertise in enterprise IT pricing.
Today is {TODAY}. Provide CURRENT approximate cost estimates for technology migrations.

CRITICAL FORMAT RULES:
- All cost values MUST be human-readable strings with dollar signs and ranges
- Example: "$2K–5K/server" or "$500–1.3K/yr" — NEVER raw numbers like 500 or 1000
- Always include the unit in the cost string (e.g. "/server", "/instance", "/yr")

For each technology, return:
- cost_upgrade: string like "$2K–5K/server (in-place upgrade)"
- cost_replace: string like "$15K–30K/server (re-platform to cloud)"
- cost_do_nothing: string like "$200–400/server/yr ESU"
- cost_unit: "server" | "instance" | "application" | "cluster" | "device"
- source_note: Brief note like "Microsoft ESU Year 2 pricing"

Reference actual vendor pricing where possible.
Return ONLY a JSON array matching the input technologies."""


def get_cost_estimates(table_data, agent=None):
    """Add cost estimate columns to each principle row. Optionally enhanced with AI."""
    # Start with defaults
    result = []
    for item in table_data:
        tech = item.get("technology", item.get("os_family", ""))
        costs = None
        for key, vals in COST_ESTIMATES.items():
            if key.lower() in tech.lower():
                costs = vals
                break
        if costs:
            item["cost_upgrade"] = costs[0]
            item["cost_replace"] = costs[1]
            item["cost_do_nothing"] = costs[2]
            item["cost_unit"] = costs[3]
        else:
            item["cost_upgrade"] = "Contact vendor"
            item["cost_replace"] = "Custom estimate required"
            item["cost_do_nothing"] = "Varies"
            item["cost_unit"] = "unit"
        item["cost_source"] = "baseline"
        result.append(item)

    # AI enhancement if agent provided
    if agent:
        try:
            tech_list = [{"technology": r.get("technology", r.get("os_family", "")),
                          "category": r.get("category", "OS"),
                          "cloud_target": r.get("cloud_target", "")}
                         for r in result]
            prompt = (
                f"Provide current cost estimates for migrating these technologies:\n"
                + json.dumps(tech_list, indent=2) +
                f"\n\nReturn JSON array with fields: technology, cost_upgrade, cost_replace, "
                f"cost_do_nothing, cost_unit, source_note"
            )
            resp = agent.client.chat.completions.create(
                model=agent.model, max_tokens=3000,
                messages=[
                    {"role": "system", "content": COST_AI_SYSTEM},
                    {"role": "user", "content": prompt}
                ])
            text = resp.choices[0].message.content.strip()
            if "```" in text:
                text = text.split("```json")[-1].split("```")[0] if "```json" in text \
                       else text.split("```")[1].split("```")[0]
            s, e = text.find("["), text.rfind("]")
            if s != -1 and e > s:
                ai_costs = json.loads(text[s:e+1])
                # Merge AI costs into results by technology name
                ai_lookup = {}
                for ac in ai_costs:
                    ai_lookup[ac.get("technology", "").lower()] = ac
                for item in result:
                    tech = item.get("technology", item.get("os_family", "")).lower()
                    if tech in ai_lookup:
                        ac = ai_lookup[tech]
                        item["cost_upgrade"] = ac.get("cost_upgrade", item["cost_upgrade"])
                        item["cost_replace"] = ac.get("cost_replace", item["cost_replace"])
                        item["cost_do_nothing"] = ac.get("cost_do_nothing", item["cost_do_nothing"])
                        item["cost_unit"] = ac.get("cost_unit", item["cost_unit"])
                        item["cost_source"] = "AI-enhanced"
                        if ac.get("source_note"):
                            item["cost_note"] = ac["source_note"]
        except Exception:
            pass  # Keep defaults if AI fails

    return result


# =============================================================================
# FEATURE: Compliance Crosswalk
# =============================================================================
COMPLIANCE_RULES = {
    "PCI DSS 4.0": {
        "description": "Payment Card Industry — requires all system components to be patched and supported",
        "eol_max_days": 0,  # No EOL software allowed
        "rule": "Req 6.3.3: All software must have active vendor security patches",
    },
    "HIPAA": {
        "description": "Health Insurance Portability — requires reasonable security safeguards",
        "eol_max_days": 90,  # 90-day grace period for patching
        "rule": "§164.312(a)(2)(iv): Encryption + patch management for ePHI systems",
    },
    "SOX (IT Controls)": {
        "description": "Sarbanes-Oxley — financial reporting system integrity",
        "eol_max_days": 180,  # 6-month remediation window
        "rule": "Section 404: IT general controls must ensure system reliability",
    },
    "NIST 800-53": {
        "description": "Federal information systems — strict patch management",
        "eol_max_days": 30,
        "rule": "SI-2: Flaw Remediation — critical patches within 30 days",
    },
    "ISO 27001": {
        "description": "Information security management — risk-based approach",
        "eol_max_days": 365,  # Risk-based; must have compensating controls
        "rule": "A.12.6.1: Technical vulnerability management within risk tolerance",
    },
}

def generate_compliance_crosswalk(wave_data):
    """Map each technology's EOL status against compliance frameworks."""
    crosswalk = []
    for item in wave_data:
        days = item.get("days_eol", 999)
        if days is None:
            days = 999
        violations = []
        warnings = []

        for framework, rules in COMPLIANCE_RULES.items():
            max_days = rules["eol_max_days"]
            if days < 0:
                # Already EOL
                violations.append({"framework": framework, "status": "VIOLATION",
                                    "detail": f"EOL {abs(int(days))}d ago — {rules['rule']}"})
            elif days <= max_days:
                violations.append({"framework": framework, "status": "AT RISK",
                                    "detail": f"EOL in {int(days)}d — within {framework} remediation window"})
            elif days <= max_days + 180:
                warnings.append({"framework": framework, "status": "MONITOR",
                                  "detail": f"EOL in {int(days)}d — approaching {framework} threshold"})

        crosswalk.append({
            **item,
            "compliance_violations": violations,
            "compliance_warnings": warnings,
            "compliance_status": "VIOLATION" if violations else ("WARNING" if warnings else "COMPLIANT"),
        })
    return crosswalk


# =============================================================================
# FEATURE: PowerPoint Export
# =============================================================================
def generate_pptx(table_data, wave_data, compliance_data, dep_data, cloud_name, selected_families):
    """Generate a PowerPoint presentation from the analysis data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(0x00, 0x1F, 0x5B)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Migration Modernization Strategy"
    p.font.size = Pt(36); p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f"Cloud Target: {cloud_name} · OS Families: {', '.join(selected_families)}"
    p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p3 = tf.add_paragraph()
    p3.text = f"Project: Apr 2026 → Jun 2028 · Generated: {datetime.now().strftime('%d %b %Y')}"
    p3.font.size = Pt(14); p3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    # ── Slide 2: Guiding Principles Table ────────────────────────────────────
    cat_order = ["OS", "Database", "Web Server", "App Server", "Framework"]
    for cat in cat_order:
        cat_rows = [r for r in table_data if r.get("category", "OS") == cat]
        if not cat_rows:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title.text_frame.paragraphs[0].text = f"{cat} — Guiding Principles"
        title.text_frame.paragraphs[0].font.size = Pt(24); title.text_frame.paragraphs[0].font.bold = True

        cols_count = 4
        rows_count = len(cat_rows) + 1
        tbl = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.1),
                                      Inches(12.3), Inches(min(rows_count * 0.6, 5.5))).table

        headers = ["Technology", "Cloud Target", "Upgrade Principle", "Replacement Principle"]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

        for i, row in enumerate(cat_rows):
            vals = [row.get("technology", ""), row.get("cloud_target", ""),
                    row.get("upgrade_principle", ""), row.get("replacement_principle", "")]
            for j, v in enumerate(vals):
                cell = tbl.cell(i + 1, j)
                cell.text = v
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                if i % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)

    # ── Slide 3: Migration Waves ─────────────────────────────────────────────
    if wave_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title.text_frame.paragraphs[0].text = "Migration Wave Plan"
        title.text_frame.paragraphs[0].font.size = Pt(24); title.text_frame.paragraphs[0].font.bold = True

        wave_colors = {1: RGBColor(0xDC,0x26,0x26), 2: RGBColor(0xEA,0x58,0x0C),
                       3: RGBColor(0xCA,0x8A,0x04), 4: RGBColor(0x16,0xA3,0x4A)}
        y = 1.2
        for w in [1, 2, 3, 4]:
            items = [r for r in wave_data if r.get("wave") == w]
            if not items:
                continue
            timeline = items[0].get("timeline", "")
            wave_name = items[0].get("wave_name", f"Wave {w}")

            box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.35))
            p = box.text_frame.paragraphs[0]
            p.text = f"{wave_name} ({timeline}) — {len(items)} technologies"
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = wave_colors.get(w, RGBColor(0,0,0))

            techs = ", ".join(f"{r.get('technology','')} ({r.get('category','')})" for r in items[:8])
            if len(items) > 8:
                techs += f" +{len(items)-8} more"
            box2 = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.35), Inches(11.5), Inches(0.3))
            box2.text_frame.paragraphs[0].text = techs
            box2.text_frame.paragraphs[0].font.size = Pt(10)
            y += 0.8

    # ── Slide 4: Compliance Summary ──────────────────────────────────────────
    if compliance_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title.text_frame.paragraphs[0].text = "Compliance Crosswalk Summary"
        title.text_frame.paragraphs[0].font.size = Pt(24); title.text_frame.paragraphs[0].font.bold = True

        violations = [r for r in compliance_data if r.get("compliance_status") == "VIOLATION"]
        warnings = [r for r in compliance_data if r.get("compliance_status") == "WARNING"]
        compliant = [r for r in compliance_data if r.get("compliance_status") == "COMPLIANT"]

        summary_text = (f"Violations: {len(violations)} technologies · "
                       f"At Risk: {len(warnings)} · Compliant: {len(compliant)}")
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
        box.text_frame.paragraphs[0].text = summary_text
        box.text_frame.paragraphs[0].font.size = Pt(16)

        if violations:
            y = 2.0
            for v in violations[:10]:
                box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.5), Inches(0.4))
                frameworks = ", ".join(cv["framework"] for cv in v.get("compliance_violations", []))
                box.text_frame.paragraphs[0].text = f"❌ {v.get('technology','')} ({v.get('category','')}) — violates: {frameworks}"
                box.text_frame.paragraphs[0].font.size = Pt(10)
                box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
                y += 0.35

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


LANDSCAPE_VERIFY_SYSTEM = """You are Agent 5 — a senior IT migration strategist.
The user says they have an OS that is not in our tracked list.
Your job is to determine:
1. Is the user referring to an OS we already track but by a different name/nickname?
   (e.g. "RHEL" = "Red Hat Enterprise Linux", "Win Server" = "Windows Server")
2. Is this a genuinely new OS that we should add to our tracking?

Our tracked OS families: Windows (10/11, Server 2003-2025), RHEL (7-10), Ubuntu (14.04-24.04),
Debian (8-13), SLES (11-16), CentOS/Stream, Rocky Linux, AlmaLinux, Oracle Linux,
macOS (Ventura/Sonoma/Sequoia), AIX (7.2-7.3), HP-UX (11i), Solaris (10-11.4),
Tru64, FreeBSD (13-14), OpenVMS, ChromeOS, Android, iOS/iPadOS, Fedora,
Raspberry Pi OS, IBM i.

Respond with ONLY this JSON:
{"match_found": true/false,
 "matched_to": "exact OS name from our list if match found, else null",
 "is_valid_os": true/false,
 "os_name": "the canonical name of the OS",
 "explanation": "brief explanation of your determination"
}"""

CONVERSATION_SYSTEM = f"""You are Agent 5 — a senior IT migration strategist at Infosys.
Help an enterprise architect define their OS and database migration policy
for a project running from 1 April 2026 to 30 June 2028.

YOUR APPROACH — BE ADAPTIVE:
- Some users will answer all questions in detail. Others will give a brief context statement
  and want to proceed quickly. BOTH are valid. Adapt to the user's pace.
- If the user gives a rich context upfront (e.g. "We are a large bank with PCI DSS scope,
  zero EOL tolerance, Azure preferred, small migration team"), extract as much as you can
  from that single message and ask only about the gaps.
- If the user says "just proceed" or "that's enough" or "go ahead" — respect that and signal ready.
- For any topics NOT discussed, apply sensible enterprise defaults in the context JSON.

TOPICS TO COVER (ask only what you still need):
1.  eol_tolerance  — Risk tolerance for EOL software (default if not asked: "Low — ESU with controls")
2.  min_runway     — Support runway at Jun 2028 (default: "12 months past project end")
3.  esu_budget     — ESU budget (default: "Limited — Tier-1 only")
4.  compliance     — Regulatory scope (default: "Internal policy only")
5.  windows_path   — Windows EOL path (default: "In-place upgrade to latest supported version")
6.  linux_path     — Linux/Unix/AIX path (default: "In-place upgrade same distro")
7.  db_eol_path    — Database EOL path (default: "In-place upgrade same vendor")
8.  oracle_stance  — Oracle licensing (default: "Evaluate case by case")
9.  cloud_provider — Cloud provider (default: "No strong preference")
10. legacy_db      — Legacy DBs stance (default: "Retain if stable, migrate if EOL in window")
11. capacity       — Migration capacity (default: "Medium — 20-30 systems in project window")
12. criticality    — System priority (default: "EOL-risk first, then compliance scope")
13. rollback       — Rollback policy (default: "Parallel run for Tier-1, in-place for others")

SIGNAL READY when:
- You have explicit or inferred answers for most topics, OR
- The user has indicated they want to proceed, OR
- You have had at least 4 exchanges and covered the most critical topics
  (eol_tolerance, compliance, windows_path, db_eol_path)

When ready, respond with ONLY this JSON:
{{"ready": true,
  "summary": "2-3 sentence org-specific policy summary based on what was discussed",
  "context": {{
    "eol_tolerance": "...", "min_runway": "...", "esu_budget": "...",
    "compliance": "...", "windows_path": "...", "linux_path": "...",
    "db_eol_path": "...", "oracle_stance": "...", "cloud_provider": "...",
    "legacy_db": "...", "capacity": "...", "criticality": "...", "rollback": "..."
  }},
  "inferred_topics": ["list of topics that used defaults, not explicitly discussed"]
}}
Today: {TODAY}. Project ends 30 Jun 2028."""

PRINCIPLES_SYSTEM = """You are a senior IT migration strategist.
Generate 8-10 Guiding Principles from a SPECIFIC policy conversation.

CRITICAL: Each principle must be ORG-SPECIFIC — use actual details from the conversation
(specific thresholds, budgets, compliance frameworks, named products, team sizes mentioned).
Generic principles like "Understand risk tolerance" are NOT acceptable.

Good example: GP-01: "PCI DSS Zero-EOL — All payment system OS/DB versions must be on
supported releases by 30 Jun 2028. No ESU permitted for Tier-1 PCI scope systems."
Bad example: GP-01: "Risk Tolerance — Understand organization's risk tolerance."

Return ONLY a JSON array:
[{"code":"GP-01","title":"4-word specific title","rule":"One specific actionable rule with named systems/dates/thresholds.",
  "trigger":"If [specific condition] → [specific action]","category":"Risk|Budget|OS|Database|Execution"}]"""

FINAL_REC_SYSTEM = f"""You are a senior IT migration strategist at Infosys.
Cross-reference:
1. Agent 2's expert technical recommendation
2. Organisation's specific policy context from the conversation
3. Agreed Guiding Principles (cite by code)

Project: 1 Apr 2026 → 30 Jun 2028. Today: {TODAY}.

For each record, produce a FINAL RECOMMENDATION that:
- Starts with: CRITICAL / UPGRADE NOW / EXTEND + PLAN / REPLACE / CLOUD MIGRATE / MONITOR
- Synthesises Agent 2's technical advice with the ORG'S SPECIFIC POLICY (not generic advice)
- Cites a GP code e.g. (GP-01)
- Is 2-3 sentences, specific to this record

Return ONLY valid JSON: {{"KEY": "VERDICT — org-specific recommendation. (GP-N)"}}"""


# ── SQLite helpers ─────────────────────────────────────────────────────────────
def _get_db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.execute("""CREATE TABLE IF NOT EXISTS conversations (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        session TEXT NOT NULL, role TEXT NOT NULL,
        content TEXT NOT NULL, ts DATETIME DEFAULT CURRENT_TIMESTAMP)""")
    conn.execute("""CREATE TABLE IF NOT EXISTS sessions (
        session TEXT PRIMARY KEY, context TEXT, summary TEXT,
        status TEXT DEFAULT 'chatting',
        created DATETIME DEFAULT CURRENT_TIMESTAMP,
        updated DATETIME DEFAULT CURRENT_TIMESTAMP)""")
    conn.commit()
    return conn

def _save_message(session_id, role, content):
    conn = _get_db()
    conn.execute("INSERT INTO conversations (session,role,content) VALUES (?,?,?)",
                 (session_id, role, content))
    conn.execute("INSERT OR IGNORE INTO sessions (session) VALUES (?)", (session_id,))
    conn.execute("UPDATE sessions SET updated=CURRENT_TIMESTAMP WHERE session=?", (session_id,))
    conn.commit(); conn.close()

def _load_messages(session_id):
    conn = _get_db()
    rows = conn.execute(
        "SELECT role,content FROM conversations WHERE session=? ORDER BY id",
        (session_id,)).fetchall()
    conn.close()
    return [{"role": r, "content": c} for r, c in rows]

def _save_session_context(session_id, context, summary, status):
    conn = _get_db()
    conn.execute("""INSERT INTO sessions (session,context,summary,status)
        VALUES (?,?,?,?) ON CONFLICT(session) DO UPDATE SET
        context=excluded.context, summary=excluded.summary,
        status=excluded.status, updated=CURRENT_TIMESTAMP""",
        (session_id, json.dumps(context), summary, status))
    conn.commit(); conn.close()

def _list_sessions():
    conn = _get_db()
    rows = conn.execute(
        "SELECT session,summary,status,updated FROM sessions ORDER BY updated DESC LIMIT 20"
    ).fetchall()
    conn.close()
    return [{"session": r[0], "summary": r[1] or "In progress",
             "status": r[2], "updated": r[3]} for r in rows]

def _delete_session(session_id):
    conn = _get_db()
    conn.execute("DELETE FROM conversations WHERE session=?", (session_id,))
    conn.execute("DELETE FROM sessions WHERE session=?", (session_id,))
    conn.commit(); conn.close()



MODEL_PREFERENCE = ["gpt-4o-mini", "gpt-3.5-turbo", "gpt-3.5-turbo-0125"]

def _resolve_model(client):
    """Try models in order, return first that works. Cache in session_state."""
    import streamlit as st
    cached = st.session_state.get("_openai_model")
    if cached:
        return cached
    for model in MODEL_PREFERENCE:
        try:
            client.chat.completions.create(
                model=model, max_tokens=5,
                messages=[{"role": "user", "content": "ping"}]
            )
            st.session_state["_openai_model"] = model
            return model
        except Exception as e:
            err = str(e).lower()
            if "not found" in err or "404" in err or "model" in err:
                continue
            st.session_state["_openai_model"] = model
            return model
    return MODEL_PREFERENCE[0]

class PolicyAnalysisAgent:

    def __init__(self, api_key: str):
        self.client = OpenAI(api_key=api_key)
        self.model  = _resolve_model(self.client)

    @staticmethod
    def get_or_create_session():
        if "a5_session_id" not in st.session_state:
            import uuid
            st.session_state.a5_session_id = str(uuid.uuid4())[:8]
        return st.session_state.a5_session_id

    @staticmethod
    def init_session():
        defaults = {
            "a5_status": "idle", "a5_context": {}, "a5_principles": [],
            "a5_costs": {}, "a5_os_done": False, "a5_db_done": False,
            "a5_ws_done": False, "a5_as_done": False, "a5_fw_done": False,
            "a5_preflight_done": False, "a5_log": [],
            "a5_landscape_families": {},
            "a5_landscape_selected": [],
            "a5_landscape_other_pending": False,
            "a5_custom_cloud_profiles": [],
        }
        for k, v in defaults.items():
            if k not in st.session_state:
                st.session_state[k] = v

    @staticmethod
    def reset():
        session_id = st.session_state.get("a5_session_id")
        if session_id:
            _delete_session(session_id)
        for k in ["a5_status","a5_session_id","a5_context","a5_principles",
                  "a5_costs","a5_os_done","a5_db_done",
                  "a5_ws_done","a5_as_done","a5_fw_done",
                  "a5_preflight_done","a5_log",
                  "a5_landscape_families","a5_landscape_selected",
                  "a5_landscape_other_pending","a5_custom_cloud_profiles",
                  "a5_principles_table_data"]:
            st.session_state.pop(k, None)
        PolicyAnalysisAgent.init_session()

    def verify_unknown_os(self, user_input: str, os_df) -> dict:
        """Check if user-described OS matches something we already track or is genuinely new."""
        known_os = ", ".join(os_df["OS Version"].unique().tolist()[:80])
        prompt = (f"The user says their environment includes: \"{user_input}\"\n\n"
                  f"Our currently tracked OS versions include (sample): {known_os}\n\n"
                  f"Is this an OS we already track (possibly by a different name), "
                  f"or is this genuinely a new OS we should add?")
        try:
            resp = self.client.chat.completions.create(
                model=self.model, max_tokens=300,
                messages=[
                    {"role": "system", "content": LANDSCAPE_VERIFY_SYSTEM},
                    {"role": "user", "content": prompt}
                ])
            text = resp.choices[0].message.content.strip()
            if "```" in text:
                text = text.split("```json")[-1].split("```")[0] if "```json" in text \
                       else text.split("```")[1].split("```")[0]
            s, e = text.find("{"), text.rfind("}")
            if s != -1 and e > s:
                return json.loads(text[s:e+1])
        except Exception:
            pass
        return {"match_found": False, "is_valid_os": False, "os_name": user_input,
                "explanation": "Could not verify — please add manually if needed."}

    def chat(self, messages: list) -> str:
        """Send conversation to gpt-4o-mini via chat.completions."""
        api_messages = [{"role": "system", "content": CONVERSATION_SYSTEM}]
        if messages:
            api_messages += messages[-20:]
        else:
            api_messages.append({"role": "user",
                                  "content": "Start the policy conversation now."})

        response = self.client.chat.completions.create(
            model=self.model,
            max_tokens=700,
            messages=api_messages
        )
        return response.choices[0].message.content.strip()

    def is_conversation_complete(self, reply: str, message_count: int = 0) -> tuple:
        """
        Accept ready signal adaptively:
        - Minimum 4 exchanges (8 messages) OR user explicitly said to proceed
        - Agent fills defaults for any uncovered topics — nothing is blocked
        Returns (done, context, summary, inferred_topics)
        """
        try:
            s, e = reply.find("{"), reply.rfind("}")
            if s != -1 and e > s:
                data = json.loads(reply[s:e+1])
                if data.get("ready"):
                    context         = data.get("context", {})
                    inferred        = data.get("inferred_topics", [])
                    summary         = data.get("summary", "")
                    # Only guard: need at least 4 exchanges (8 messages)
                    # so the agent has had a chance to ask something
                    if message_count < 8:
                        return False, {}, "", []
                    # Fill any completely missing keys with defaults
                    DEFAULTS = {
                        "eol_tolerance":  "Low — ESU with compensating controls",
                        "min_runway":     "12 months past Jun 2028",
                        "esu_budget":     "Limited — Tier-1 critical systems only",
                        "compliance":     "Internal policy only",
                        "windows_path":   "In-place upgrade to latest supported version",
                        "linux_path":     "In-place upgrade same distro",
                        "db_eol_path":    "In-place upgrade same vendor",
                        "oracle_stance":  "Evaluate Oracle vs alternatives case by case",
                        "cloud_provider": "No strong cloud preference",
                        "legacy_db":      "Retain if stable, migrate if EOL within project",
                        "capacity":       "Medium — 20-30 systems across project lifespan",
                        "criticality":    "EOL-risk first, then compliance scope",
                        "rollback":       "Parallel run for Tier-1, in-place for others",
                    }
                    for k, default in DEFAULTS.items():
                        if not context.get(k, "").strip():
                            context[k] = f"[Default] {default}"
                            if k not in inferred:
                                inferred.append(k)
                    return True, context, summary, inferred
        except Exception:
            pass
        return False, {}, "", []

    def generate_principles(self, context: dict, session_id: str) -> list:
        messages = _load_messages(session_id)
        conv = "\n".join(
            f"{'ARCHITECT' if m['role']=='user' else 'AGENT5'}: {m['content'][:200]}"
            for m in messages[-24:])
        ctx  = "\n".join(f"{k}: {v}" for k, v in context.items())
        try:
            resp = self.client.chat.completions.create(
                model=self.model, max_tokens=2500,
                messages=[
                    {"role": "system", "content": PRINCIPLES_SYSTEM},
                    {"role": "user",   "content": f"CONTEXT:\n{ctx}\n\nCONVERSATION:\n{conv}"}
                ])
            text = resp.choices[0].message.content.strip()
            if "```" in text:
                text = text.split("```json")[-1].split("```")[0] if "```json" in text \
                       else text.split("```")[1].split("```")[0]
            s, e = text.find("["), text.rfind("]")
            if s != -1 and e > s:
                return json.loads(text[s:e+1])
        except Exception:
            pass
        return [{"code": "GP-01", "title": "Zero EOL Tolerance",
                 "rule": "No EOL software past 30 Jun 2028.",
                 "trigger": "EOL before Jun 2028 → Upgrade", "category": "Risk"}]

    def fetch_costs(self, progress_cb=None) -> dict:
        queries = [
            ("Windows Server ESU",
             "Windows Server 2016 and 2019 Extended Security Updates estimated cost per server per year. Include Azure Arc free ESU option."),
            ("SQL Server ESU",
             "SQL Server 2016 2017 2019 Extended Security Updates estimated cost per core per year."),
            ("Oracle Database Support",
             "Oracle Database annual support cost approximately 22 percent of license. Extended support surcharge details."),
            ("RHEL Subscription",
             "Red Hat Enterprise Linux RHEL Standard and Premium subscription cost per server per year estimates."),
            ("Cloud DB Pricing",
             "Azure SQL Managed Instance and AWS RDS Aurora PostgreSQL approximate monthly cost per instance."),
            ("Migration Tools",
             "AWS Database Migration Service and Azure Database Migration Service approximate cost per hour."),
        ]
        costs = {}
        fallbacks = {
            "Windows Server ESU":
                "Windows Server 2016 ESU: ~$198/server/yr (Yr1), ~$396 (Yr2). Free via Azure Arc. Server 2019: ESU from 2029.",
            "SQL Server ESU":
                "SQL Server 2016 ESU: ~$1,418/core/yr. SQL Server 2019 mainstream support ends Jan 2025; extended until Jan 2030 (no ESU needed yet).",
            "Oracle Database Support":
                "Oracle annual support: ~22% of license list price. Extended Support (years 4-5) adds 10% surcharge (~32% total).",
            "RHEL Subscription":
                "RHEL Standard: ~$800/yr/socket-pair. RHEL Premium (24x7 support): ~$1,300/yr/socket-pair.",
            "Cloud DB Pricing":
                "Azure SQL MI (4 vCores, GP): ~$465/mo. AWS RDS PostgreSQL db.t3.medium: ~$60/mo. Aurora PostgreSQL: ~$200/mo.",
            "Migration Tools":
                "AWS DMS: ~$0.18/hr per replication instance + data transfer. Azure DMS Standard: ~$0.10/hr.",
        }
        for i, (vendor, query) in enumerate(queries):
            if progress_cb:
                progress_cb(i / len(queries), f"Fetching cost estimate: {vendor}...")
            try:
                resp = self.client.chat.completions.create(
                    model=self.model, max_tokens=300,
                    messages=[{"role": "user", "content":
                        f"Based on your training knowledge, give a 2-3 sentence cost estimate for: {query}\n"
                        f"Note that figures are approximate and may have changed since your training cutoff."}])
                result = resp.choices[0].message.content.strip()
                costs[vendor] = result or fallbacks[vendor]
            except Exception:
                costs[vendor] = fallbacks[vendor]
        if progress_cb:
            progress_cb(1.0, "✅ Cost estimates ready.")
        return costs

    def analyse_os(self, df, principles, costs, context, progress_cb=None):
        return self._analyse(df, "OS", principles, costs, context, progress_cb)

    def analyse_db(self, df, principles, costs, context, progress_cb=None):
        return self._analyse(df, "DB", principles, costs, context, progress_cb)

    def analyse_ws(self, df, principles, costs, context, progress_cb=None):
        return self._analyse(df, "WS", principles, costs, context, progress_cb)

    def analyse_as(self, df, principles, costs, context, progress_cb=None):
        return self._analyse(df, "AS", principles, costs, context, progress_cb)

    def analyse_fw(self, df, principles, costs, context, progress_cb=None):
        return self._analyse(df, "FW", principles, costs, context, progress_cb)

    def _analyse(self, df, kind, principles, costs, context, progress_cb):
        df = df.copy()
        for col in ["Final Recommendation", "Final Verdict", "Analysis Source"]:
            if col not in df.columns:
                df[col] = ""

        gp_text   = "\n".join(f"{p['code']}: {p['title']} — {p['rule']}" for p in principles)
        cost_text = "\n".join(f"{v}: {s}" for v, s in costs.items())
        ctx_text  = "\n".join(f"{k}: {v}" for k, v in context.items())
        rows = df.to_dict("records")
        total = len(rows); batch = 15
        ai_count = 0; rb_count = 0

        for i in range(0, total, batch):
            chunk = rows[i:i+batch]
            if progress_cb:
                progress_cb(i / total,
                    f"🧠 Generating Final Recommendations — {kind} rows {i+1}–{min(i+batch,total)} of {total}...")

            if kind == "OS":
                rows_text = "\n".join(
                    f"KEY={r['OS Version']} | Mainstream={r.get('Mainstream/Full Support End','')} | "
                    f"Extended={r.get('Extended/LTSC Support End','')} | "
                    f"Agent2={r.get('Recommendation','')[:120]}"
                    for r in chunk
                )
            else:
                # DB, WS, AS, FW all share same column structure
                name_col_map = {"DB": "Database", "WS": "Web Server", "AS": "App Server", "FW": "Framework"}
                name_col = name_col_map.get(kind, "Database")
                rows_text = "\n".join(
                    f"KEY={r.get(name_col,'?')} {r.get('Version','?')} | Status={r.get('Status','')} | "
                    f"Extended={r.get('Extended Support End','')} | "
                    f"Replace={r.get('Replace','')} | Alt={r.get('Primary Alternative','')} | "
                    f"Agent2={r.get('Recommendation','')[:120]}"
                    for r in chunk
                )

            prompt = (f"ORG POLICY:\n{ctx_text}\n\nGUIDING PRINCIPLES:\n{gp_text}\n\n"
                      f"VENDOR COSTS:\n{cost_text}\n\n"
                      f"PROJECT: 1 Apr 2026 → 30 Jun 2028 | Today: {TODAY}\n\n"
                      f"RECORDS:\n{rows_text}\n\n"
                      f"Return ONLY JSON: {{\"KEY\": \"VERDICT — recommendation. (GP-N)\"}}")

            recs = {}; api_worked = False
            for attempt in range(2):
                try:
                    import time
                    if attempt > 0: time.sleep(2)
                    resp = self.client.chat.completions.create(
                        model=self.model, max_tokens=4000,
                        messages=[
                            {"role": "system", "content": FINAL_REC_SYSTEM},
                            {"role": "user",   "content": prompt}
                        ])
                    text = resp.choices[0].message.content.strip()
                    if "```" in text:
                        text = text.split("```json")[-1].split("```")[0] if "```json" in text \
                               else text.split("```")[1].split("```")[0]
                    s, e = text.find("{"), text.rfind("}")
                    parsed = json.loads(text[s:e+1]) if s != -1 and e > s else {}
                    if parsed:
                        recs = parsed; api_worked = True; break
                except Exception as ex:
                    last_error = str(ex)

            for j, row in enumerate(chunk):
                if kind == "OS":
                    key = row["OS Version"]
                else:
                    nc = {"DB": "Database", "WS": "Web Server", "AS": "App Server", "FW": "Framework"}.get(kind, "Database")
                    key = f"{row.get(nc, '?')} {row.get('Version', '')}".strip()
                if key in recs:
                    rec = recs[key]; source = "OpenAI"; ai_count += 1
                else:
                    rec = self._rule_based(row, kind); source = "Rule-based"; rb_count += 1
                verdict = next((v for v in VERDICTS if rec.upper().startswith(v)), "MONITOR")
                df.at[i+j, "Final Recommendation"] = rec
                df.at[i+j, "Final Verdict"]        = verdict
                df.at[i+j, "Analysis Source"]      = source

            import time as _t; _t.sleep(0.3)

        if progress_cb:
            progress_cb(1.0, f"✅ {kind} done — OpenAI: {ai_count} | Rule-based: {rb_count}")
        return df

    def _rule_based(self, row, kind):
        def _parse(s):
            try: return datetime.strptime(str(s)[:10], "%Y-%m-%d").date()
            except: return None
        if kind == "OS":
            end  = _parse(row.get("Extended/LTSC Support End","")) \
                or _parse(row.get("Mainstream/Full Support End",""))
            name = row.get("OS Version", "?")
        else:
            # DB, WS, AS, FW all share same column structure
            end  = _parse(row.get("Extended Support End","")) \
                or _parse(row.get("Mainstream / Premier End",""))
            nc = {"DB": "Database", "WS": "Web Server", "AS": "App Server", "FW": "Framework"}.get(kind, "Database")
            name = f"{row.get(nc,'?')} {row.get('Version','?')}"
        if not end:
            return f"MONITOR — Verify support dates for {name}. (GP-02)"
        if end < TODAY:
            return f"CRITICAL — {name} EOL as of {end}. Immediate action required. (GP-01)"
        if end < PROJECT_END:
            return f"UPGRADE NOW — {name} ends {end}, before Jun 2028. Act now. (GP-01)"
        if end < date(2030, 6, 30):
            return f"EXTEND + PLAN — {name} until {end}. Plan upgrade within project. (GP-02)"
        return f"MONITOR — {name} until {end}, past project end. No immediate action. (GP-05)"
